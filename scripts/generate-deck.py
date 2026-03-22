"""
Empire 8 Sales Direct — Branded Sales Deck Generator v5
Interior template centered on slide, text inside gold frame area.
Fixed header/body overlap. Expanded content. Added SMS pipeline.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
HEADER_IMG = os.path.join(os.path.expanduser("~"), "Desktop", "empire8-header-splash.png")
INTERIOR_IMG = os.path.join(os.path.expanduser("~"), "Desktop", "empire8-interior-base.png")
LOGO_IMG = os.path.join(BASE, "public", "logo.png")
ZONE_MAP = os.path.join(BASE, "public", "zone-map.jpg")
OUT_PATH = os.path.join(os.path.expanduser("~"), "Desktop", "Empire8_Sales_Deck_v5c.pptx")

# Colors
PURPLE_DARK = RGBColor(0x0F, 0x05, 0x20)
PURPLE_MID = RGBColor(0x2D, 0x0A, 0x4E)
PURPLE_ROYAL = RGBColor(0x4A, 0x0E, 0x78)
GOLD = RGBColor(0xC8, 0xA2, 0x3C)
GOLD_LIGHT = RGBColor(0xE8, 0xD4, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_70 = RGBColor(0xB3, 0xB3, 0xB3)
WHITE_50 = RGBColor(0x80, 0x80, 0x80)
TEXT_DARK = RGBColor(0x1A, 0x10, 0x25)
TEXT_MED = RGBColor(0x3A, 0x3A, 0x3A)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)
BG_LAVENDER = RGBColor(0xF0, 0xE6, 0xF8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# Interior image is 850x1100 (portrait). On a 16:9 slide (13.333 x 7.5):
INT_H = SLIDE_H
INT_W = Inches(7.5 * (850.0 / 1100.0))  # ~5.795"
INT_X = Emu(int((SLIDE_W - INT_W) / 2))  # centered

# Gold frame boundaries (translated to slide coordinates)
FRAME_LEFT = Inches(3.769 + 5.795 * 0.12)   # ~4.46"
FRAME_TOP = Inches(7.5 * 0.22)              # ~1.65"
FRAME_WIDTH = Inches(5.795 * 0.76)          # ~4.40"
FRAME_BOTTOM = Inches(7.5 * 0.88)          # ~6.60"
FRAME_HEIGHT = FRAME_BOTTOM - FRAME_TOP    # ~4.95"


def add_bg(slide, color=PURPLE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_interior_slide(slide):
    """Lavender background + centered interior template image."""
    add_bg(slide, BG_LAVENDER)
    slide.shapes.add_picture(INTERIOR_IMG, INT_X, Emu(0), INT_W, INT_H)


def tb(slide, left, top, width, height, text, size=14,
       color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT,
       font="Calibri", lsp=1.4, italic=False):
    """Shorthand text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(size * lsp)
    return tf


def bullets(slide, left, top, width, height, items,
            size=13, color=TEXT_MED, lsp=1.6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.line_spacing = Pt(size * lsp)
        p.space_after = Pt(3)
    return tf


def section_label(slide, text):
    """Purple banner with gold text, positioned at top of gold frame."""
    label_y = FRAME_TOP + Inches(0.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    FRAME_LEFT, label_y,
                                    FRAME_WIDTH, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE_MID
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {text.upper()}"
    p.font.size = Pt(11)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = "Calibri"


def slide_title(slide, text):
    """Title positioned below section label with clear gap."""
    title_y = FRAME_TOP + Inches(0.55)
    tb(slide, FRAME_LEFT + Inches(0.15), title_y,
       FRAME_WIDTH - Inches(0.3), Inches(0.45),
       text, size=18, color=TEXT_DARK, bold=True)


def footer(slide):
    tb(slide, FRAME_LEFT, FRAME_BOTTOM - Inches(0.35),
       FRAME_WIDTH, Inches(0.3),
       "mikesutherland@empire8.com  \u2022  empire8ny.com",
       size=7, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# Content area inside the frame — INCREASED GAP to prevent header overlap
def CL():  # content left
    return FRAME_LEFT + Inches(0.2)

def CT():  # content top (below title) — pushed down to 1.5 to fully clear title
    return FRAME_TOP + Inches(1.5)

def CWF():  # content width inside frame
    return FRAME_WIDTH - Inches(0.4)

def CH():  # content height available (adjusted for new CT)
    return FRAME_HEIGHT - Inches(1.85)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PURPLE_DARK)

slide.shapes.add_picture(HEADER_IMG, Emu(0), Emu(0), SLIDE_W, Inches(4.0))

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Inches(3.2), SLIDE_W, Inches(4.3))
shape.fill.solid()
shape.fill.fore_color.rgb = PURPLE_DARK
shape.line.fill.background()

logo_h = Inches(5.5)
logo_w = Inches(5.5 * (853.0 / 1280.0))
logo_x = SLIDE_W - logo_w - Inches(0.8)
slide.shapes.add_picture(LOGO_IMG, logo_x, Inches(0.3), logo_w, logo_h)

tb(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.2),
   "WE MAKE CANNABIS\nPRODUCTS SELL.",
   size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER, lsp=1.1)

tb(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
   "Brand Ambassadors  \u2022  AI Marketing  \u2022  Retail Intelligence  \u2022  SMS Outreach",
   size=15, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

tb(slide, Inches(2), Inches(6.0), Inches(9.3), Inches(0.4),
   "Empire 8 Sales Direct  \u2014  New York\u2019s Premier Cannabis Sales Partner",
   size=11, color=WHITE_50, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Executive Summary")
slide_title(slide, "Who We Are")

body = (
    "Empire 8 is a wholesale-direct company combining manufacturing, "
    "high-speed logistics, and a proprietary AI-powered sales and marketing platform.\n\n"
    "We combine vertical integration without compromising product quality, "
    "speed-to-market, or brand integrity. No bloated commissions. No middlemen.\n\n"
    "We are dedicated to creating lasting relationships between quality brands "
    "and responsible dispensary operators. We will NOT accept the easy sale to "
    "inflate your A/R with bad actors who have no intention to pay.\n\n"
    "Our four pillars: Wholesale direct pricing \u2022 AI-powered marketing \u2022 "
    "In-store ambassadors \u2022 Personalized SMS outreach"
)
tb(slide, CL(), CT(), CWF(), CH(),
   body, size=9, color=TEXT_MED, lsp=1.4)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.35),
   "We have rolled our sleeves up \u2014 come see what an honest day\u2019s work looks like.",
   size=10, color=PURPLE_ROYAL, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: THE PROBLEM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "The Problem")
slide_title(slide, "Why Products Fail on the Shelf")

problems = [
    "Products sit on shelves collecting dust \u2014 no one is actively selling them",
    "Budtenders push what they know, not what\u2019s best for the customer",
    "Brands spend thousands on marketing that never converts to in-store sales",
    "No visibility into what happens at store level after delivery",
    "Dispensary operators don\u2019t have time to learn every SKU",
    "Distributors drop product and disappear \u2014 no follow-up, no sell-through",
    "No direct communication channel to the end consumer",
    "Reorders stall \u2014 no demand creation happening between deliveries",
]
bullets(slide, CL(), CT(), CWF(), CH(), problems, size=9, lsp=1.45)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: THE SHIFT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "The Shift")
slide_title(slide, "Attention + Influence = Sell-Through")

shifts = [
    "Customers discover products on social media before entering a dispensary",
    "Budtenders influence 60\u201380% of all cannabis purchases at point of sale",
    "Social proof drives trust \u2014 reviews, reels, budtender picks, customer stories",
    "The brands that win are the ones customers ask for by name",
    "Personalized SMS from the dispensary keeps customers coming back for specific products",
    "Demand must be created before the sale \u2014 not after the product is already on the shelf",
]
bullets(slide, CL(), CT(), CWF(), CH(), shifts, size=10, lsp=1.65)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "The old model is broken. This is the new playbook.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: 3-PILLAR SYSTEM → NOW 4 PILLARS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "What Makes Empire 8 Different")
slide_title(slide, "The 4-Pillar System")

pillars = [
    ("01  In-Store Influence",
     "Trained brand ambassadors deployed to your territory. They build relationships "
     "with budtenders, run demos, and make sure your products get recommended first."),
    ("02  Retail Intelligence",
     "Data-driven deployment. We track what sells, where it sells, and why. "
     "Ambassadors are deployed based on real performance data \u2014 not guesswork."),
    ("03  AI Marketing Engine",
     "Every store visit becomes content. Product reels, budtender picks, behind-the-scenes. "
     "Customers walk in already knowing what to buy."),
    ("04  Personalized SMS Pipeline",
     "AI-generated personalized text messages sent directly to dispensary clientele. "
     "Targeted promotions, new drop alerts, and re-engagement \u2014 driving repeat foot traffic."),
]

for i, (title, desc) in enumerate(pillars):
    y = CT() + Inches(i * 0.82)
    tb(slide, CL(), y, CWF(), Inches(0.22),
       title, size=11, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL(), y + Inches(0.24), CWF(), Inches(0.5),
       desc, size=9, color=TEXT_MED, lsp=1.35)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: PILLAR 1 DEEP DIVE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Pillar 1")
slide_title(slide, "In-Store Influence")

tb(slide, CL(), CT(), CWF(), Inches(0.35),
   "Our ambassadors are a sales force and media team combined. They don\u2019t just "
   "drop off product \u2014 they actively sell it.",
   size=10, color=TEXT_MED, lsp=1.4)

items = [
    "Educate budtenders on features, effects, strain profiles, and talking points",
    "Build genuine relationships so budtenders recommend your product over competitors",
    "Capture real-time data \u2014 shelf placement, competitor products, inventory levels",
    "Create authentic content during every store visit (reels, photos, budtender testimonials)",
    "Run product demos and sampling events that drive immediate sell-through",
    "Deliver weekly visit reports with actionable insights and reorder recommendations",
    "Identify underperforming locations and recommend corrective actions",
]
bullets(slide, CL(), CT() + Inches(0.4), CWF(), CH() - Inches(0.4),
        items, size=9, lsp=1.55)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: AMBASSADOR MODEL
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Ambassador Model")
slide_title(slide, "How Our Ambassadors Work")

steps = [
    ("1. DEPLOY", "Assigned to territories based on dispensary density, brand priority, and market opportunity. "
     "Each ambassador covers 15\u201325 dispensaries on a rotating weekly schedule."),
    ("2. EDUCATE", "Deep product knowledge sessions with budtenders \u2014 effects, comparisons, customer profiles. "
     "Ambassadors carry branded leave-behinds and point-of-sale materials."),
    ("3. INFLUENCE", "Build genuine relationships so budtenders recommend your product first. "
     "Incentive programs, sampling, and consistent face time create loyalty."),
    ("4. CAPTURE", "Shelf photos, competitor intel, inventory counts, budtender feedback, and customer reactions \u2014 "
     "all logged into our CRM in real-time during every visit."),
    ("5. CREATE", "Social content from every visit \u2014 product reels, budtender picks, store highlights, "
     "behind-the-scenes footage. Content is edited and published within 24 hours."),
    ("6. REPORT", "Weekly performance reports: stores visited, conversations had, products moved, "
     "reorder recommendations, and competitive positioning updates."),
]

for i, (title, desc) in enumerate(steps):
    y = CT() + Inches(i * 0.52)
    tb(slide, CL(), y, Inches(1.0), Inches(0.2),
       title, size=9, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL() + Inches(1.05), y, CWF() - Inches(1.05), Inches(0.45),
       desc, size=8, color=TEXT_MED, lsp=1.3)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: RETAIL INTELLIGENCE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Pillar 2")
slide_title(slide, "Retail Intelligence")

tb(slide, CL(), CT(), CWF(), Inches(0.3),
   "We deploy based on real performance metrics, not guesswork. We track:",
   size=9, color=TEXT_MED, lsp=1.3)

tracked = [
    "Budtender knowledge \u2014 can they recommend your product confidently?",
    "Store engagement \u2014 receptive, reordering, displaying properly, or cold?",
    "Shelf visibility \u2014 placement, positioning, proximity to competitors",
    "Competitive landscape \u2014 what\u2019s next to you and at what price?",
    "Inventory levels \u2014 real-time stock counts, out-of-stock alerts",
    "Sales velocity \u2014 units moved per week per location",
    "Customer feedback \u2014 what buyers say at the register",
]
bullets(slide, CL(), CT() + Inches(0.3), CWF(), CH() - Inches(0.3),
        tracked, size=9, lsp=1.4)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Data drives every deployment decision.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: AI MARKETING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Pillar 3")
slide_title(slide, "AI Marketing Engine")

content = [
    "Every store visit becomes content \u2014 AI edits raw footage into polished posts within 24 hours",
    "Product reels and highlight videos with hooks and captions",
    "Budtender picks \u2014 authentic on-camera recommendations",
    "Behind-the-scenes content showing people and process",
    "Store spotlight features \u2014 cross-promote dispensary partners",
    "Brand storytelling \u2014 origin stories, product deep-dives",
    "Geo-targeted campaigns \u2014 dispensary tagging, local hashtags",
    "AI-written captions optimized for each platform\u2019s algorithm",
]
bullets(slide, CL(), CT(), CWF(), CH(),
        content, size=9, lsp=1.35)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Customers walk in already knowing what to buy.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: PERSONALIZED SMS PIPELINE (NEW)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Pillar 4")
slide_title(slide, "Personalized SMS Pipeline")

tb(slide, CL(), CT(), CWF(), Inches(0.3),
   "AI-generated personalized texts sent to dispensary customers. "
   "Tailored to purchase history, preferences, and local inventory.",
   size=9, color=TEXT_MED, lsp=1.3)

sms_items = [
    "AI generates unique messages for each customer segment",
    "New drop alerts \u2014 \u201cHey [Name], that strain you loved is back at [Store]\u201d",
    "Targeted promotions \u2014 discounts and bundles based on past purchases",
    "Re-engagement \u2014 win back lapsed customers with personalized outreach",
    "Dispensary-branded \u2014 messages come from the store, not Empire 8",
    "Compliance-first \u2014 opt-in/opt-out managed, within state regulations",
    "Performance tracking \u2014 open rates, redemptions, foot traffic attribution",
]
bullets(slide, CL(), CT() + Inches(0.3), CWF(), CH() - Inches(0.3),
        sms_items, size=9, lsp=1.4)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Turn one-time buyers into repeat customers on autopilot.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 11: SMS PIPELINE — HOW IT WORKS (NEW)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "SMS Pipeline")
slide_title(slide, "How the Pipeline Works")

steps_sms = [
    ("1. DATA INTAKE",
     "Dispensary shares their customer list (with consent). We enrich it with "
     "purchase history, product preferences, and visit frequency."),
    ("2. AI SEGMENTATION",
     "Our AI clusters customers into segments: loyal regulars, lapsed buyers, "
     "deal hunters, strain-specific enthusiasts, new customers, and high-value VIPs."),
    ("3. MESSAGE GENERATION",
     "AI writes personalized messages for each segment. Every text feels like it came "
     "from a budtender who knows you \u2014 not a mass blast."),
    ("4. DISPENSARY APPROVAL",
     "Messages are reviewed and approved by the dispensary before sending. "
     "Their brand, their voice, their customer relationship."),
    ("5. SEND + TRACK",
     "Messages go out via compliant SMS infrastructure. We track opens, "
     "click-throughs, store visits, and actual purchases attributed to each campaign."),
    ("6. OPTIMIZE",
     "AI learns what works. Message timing, tone, offers, and segments are "
     "continuously refined based on real conversion data."),
]

for i, (title, desc) in enumerate(steps_sms):
    y = CT() + Inches(i * 0.5)
    tb(slide, CL(), y, Inches(1.2), Inches(0.2),
       title, size=9, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL() + Inches(1.25), y, CWF() - Inches(1.25), Inches(0.42),
       desc, size=8, color=TEXT_MED, lsp=1.3)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: SOCIAL ADVANTAGE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "The Social Advantage")
slide_title(slide, "Everyone Wins")

groups = [
    ("Dispensaries:",
     "Free promotion, increased foot traffic, professional online presence, "
     "personalized SMS campaigns that drive repeat visits, and a dedicated "
     "brand ambassador who becomes part of their team."),
    ("Brands:",
     "Visibility in stores AND online simultaneously. Your product gets recommended "
     "by trained budtenders, featured in social content, and promoted directly to "
     "consumers via SMS \u2014 a complete sell-through system."),
    ("Customers:",
     "Discover products before they shop through social media. Get personalized "
     "recommendations from budtenders who actually know the product. Receive "
     "text alerts about new drops and deals they actually care about."),
]

for i, (title, desc) in enumerate(groups):
    y = CT() + Inches(i * 0.95)
    tb(slide, CL(), y, CWF(), Inches(0.22),
       title, size=11, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL(), y + Inches(0.25), CWF(), Inches(0.6),
       desc, size=9, color=TEXT_MED, lsp=1.35)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Demand is created before the sale.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 13: TERRITORY DEVELOPMENT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Territory Development")
slide_title(slide, "Phase 1 \u2014 Strategic Coverage")

phase1 = [
    "NYC Metro \u2014 5 ambassadors (Manhattan, Brooklyn, Queens, Bronx)",
    "Long Island \u2014 1 ambassador (Nassau + Suffolk County)",
    "Capital Region \u2014 2 ambassadors (Albany, Troy, Saratoga)",
    "Rochester Metro \u2014 2 ambassadors (Monroe County + surrounding)",
    "Buffalo / WNY \u2014 2 ambassadors (Erie County + Niagara)",
]
bullets(slide, CL(), CT(), CWF(), Inches(1.8),
        phase1, size=9, lsp=1.5)

tb(slide, CL(), CT() + Inches(1.9), CWF(), Inches(0.3),
   "Total Phase 1: 12 Ambassadors \u2022 ~200 Dispensaries",
   size=11, color=PURPLE_ROYAL, bold=True)

tb(slide, CL(), CT() + Inches(2.2), CWF(), Inches(0.25),
   "Each ambassador visits 15\u201325 stores per week on a rotating schedule.",
   size=9, color=TEXT_MED, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 14: COVERAGE STRATEGY
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Coverage Strategy")
slide_title(slide, "Lean Coverage. Maximum Influence.")

strategies = [
    "Focus on top dispensaries \u2014 highest-volume, highest-influence stores first",
    "Build deep relationships \u2014 consistent, in-person presence every week creates trust",
    "Performance-based expansion \u2014 growth is earned with data, never rushed on hope",
    "Ambassador tenure matters \u2014 same face, same store, real relationships over time",
    "Competitive exclusivity \u2014 we won\u2019t represent competing brands in the same category",
]
bullets(slide, CL(), CT(), CWF(), Inches(2.5),
        strategies, size=10, lsp=1.7)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Depth over breadth. Quality over quantity.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 15: SERVICES OVERVIEW
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Services Overview")
slide_title(slide, "Three Growth Engines")

services = [
    ("1. Ambassador Services",
     "In-store brand ambassadors who educate budtenders, build relationships, "
     "capture competitive intelligence, and actively sell your product at the point of sale."),
    ("2. AI Marketing Services",
     "Content creation, social media management, dispensary tagging, geo-targeted campaigns, "
     "and brand storytelling \u2014 all AI-powered for speed and scale."),
    ("3. Personalized SMS Outreach",
     "AI-generated personalized text messages to dispensary customer lists. "
     "New drop alerts, targeted promotions, and re-engagement campaigns that drive repeat foot traffic."),
]

for i, (title, desc) in enumerate(services):
    y = CT() + Inches(i * 0.9)
    tb(slide, CL(), y, CWF(), Inches(0.22),
       title, size=11, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL(), y + Inches(0.25), CWF(), Inches(0.55),
       desc, size=9, color=TEXT_MED, lsp=1.35)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Individually powerful. Together unstoppable.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 16: MARKETING PACKAGES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Marketing Packages")
slide_title(slide, "Choose Your Growth Plan")

tiers_text = (
    "SMALL \u2014 $2,500/month\n"
    "\u2022  12\u201316 posts/month across 1 platform\n"
    "\u2022  AI-written captions with hooks and CTAs\n"
    "\u2022  Monthly content calendar\n"
    "\u2022  Basic performance reporting\n\n"
    "FLEX \u2014 $5,000/month  (Most Popular)\n"
    "\u2022  20\u201330 posts/month across multiple platforms\n"
    "\u2022  Product launch campaigns + dispensary tagging\n"
    "\u2022  Budtender pick content series\n"
    "\u2022  Bi-weekly performance insights + optimization\n\n"
    "BOSS \u2014 $9,000/month  (Full Service)\n"
    "\u2022  40\u201360 posts/month with daily strategy\n"
    "\u2022  Influencer coordination + brand partnerships\n"
    "\u2022  Geo-targeted ad campaigns\n"
    "\u2022  Weekly optimization reports\n"
    "\u2022  SMS pipeline integration included"
)
tb(slide, CL(), CT(), CWF(), CH(),
   tiers_text, size=9, color=TEXT_MED, lsp=1.3)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 17: SERVICE COMPARISON
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Service Comparison")
slide_title(slide, "Choose Your Growth Path")

rows = [
    ("Marketing Only", "Customers discover your product online and ask for it by name"),
    ("Ambassadors Only", "Budtenders actively push your product at the point of sale"),
    ("SMS Only", "Repeat customers come back specifically for your product"),
    ("Full System", "Maximum sell-through \u2014 demand created before, during, and after the sale"),
]

for i, (svc, outcome) in enumerate(rows):
    y = CT() + Inches(i * 0.6)
    tb(slide, CL(), y, Inches(1.3), Inches(0.22),
       svc, size=10, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL() + Inches(1.35), y, CWF() - Inches(1.35), Inches(0.45),
       outcome, size=10, color=TEXT_MED, lsp=1.35)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Best results come from combining all three.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 18: AMBASSADOR PRICING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Pricing (Ambassadors)")
slide_title(slide, "In-Store Influence Packages")

pricing_items = [
    ("Starter \u2014 $4,000/month",
     "1 ambassador, 1 territory, 15\u201320 stores per week. Weekly reporting. "
     "Ideal for brands entering a new market."),
    ("Growth \u2014 $6,500/month",
     "1\u20132 ambassadors, expanded territory, 25\u201340 stores. Content creation included. "
     "Bi-weekly strategy calls."),
    ("Scale \u2014 $9,000\u2013$12,000/month",
     "2\u20133 ambassadors, full regional coverage, 40\u201360+ stores. Priority scheduling, "
     "demo events, and dedicated account manager."),
    ("\u00c0 la carte sessions \u2014 $450\u2013$750/session",
     "Single-day activations, product launches, sampling events, or store openings."),
]

for i, (tier, desc) in enumerate(pricing_items):
    y = CT() + Inches(i * 0.7)
    tb(slide, CL(), y, CWF(), Inches(0.22),
       tier, size=10, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL(), y + Inches(0.24), CWF(), Inches(0.4),
       desc, size=9, color=TEXT_MED, lsp=1.3)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 19: SMS PRICING (NEW)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Pricing (SMS Pipeline)")
slide_title(slide, "Personalized Outreach Packages")

sms_pricing = [
    ("Starter \u2014 $1,500/month",
     "Up to 2,000 personalized messages/month. 2 campaigns. "
     "Basic segmentation. Monthly performance report."),
    ("Growth \u2014 $3,000/month",
     "Up to 5,000 messages/month. 4 campaigns. Advanced AI segmentation. "
     "A/B testing. Bi-weekly optimization."),
    ("Scale \u2014 $5,000/month",
     "Unlimited messages. Weekly campaigns. Full AI optimization loop. "
     "Custom integrations with dispensary POS systems. Dedicated strategist."),
    ("Add-on with Ambassador or Marketing package \u2014 20% discount",
     "Bundle SMS with any existing Empire 8 service for maximum sell-through."),
]

for i, (tier, desc) in enumerate(sms_pricing):
    y = CT() + Inches(i * 0.7)
    tb(slide, CL(), y, CWF(), Inches(0.22),
       tier, size=10, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL(), y + Inches(0.24), CWF(), Inches(0.4),
       desc, size=9, color=TEXT_MED, lsp=1.3)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 20: UNIT ECONOMICS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Unit Economics")
slide_title(slide, "Strong Cash Flow at Every Level")

economics = [
    "Ambassador cost to Empire 8: ~$2K/month per ambassador",
    "Ambassador service margin: 45\u201360% depending on package tier",
    "Single session margins: 75%+ on \u00e0 la carte activations",
    "Marketing services: 60\u201370% margin (AI reduces production cost dramatically)",
    "SMS pipeline: 80%+ margin (AI generation, bulk SMS rates, minimal manual labor)",
    "Blended margin across all services: 55\u201370% at scale",
    "Revenue compounds: each new brand multiplies ambassador utilization without linear cost increase",
]
bullets(slide, CL(), CT(), CWF(), CH(),
        economics, size=9, lsp=1.55)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Capital-light model. Cash flow positive from month one.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 21: GROWTH ROADMAP — SCALING THE NETWORK (EXPANDED)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Growth Road Map")
slide_title(slide, "Scaling the Network")

phases = [
    ("Phase 1 \u2014 Foundation (Months 1\u20133)",
     "6\u201310 founding brand partners. 12 ambassadors across 5 NY territories. "
     "Prove the model, generate case studies, refine the playbook. "
     "Launch SMS pipeline with 3\u20135 dispensary partners."),
    ("Phase 2 \u2014 Density (Months 4\u20136)",
     "Expand to 15\u201320 brands. Add 6\u20138 ambassadors to deepen NYC and Long Island coverage. "
     "Onboard 15+ dispensaries onto the SMS pipeline. Launch AI marketing for 10+ brands. "
     "Shared ambassador model reduces per-brand cost as portfolio grows."),
    ("Phase 3 \u2014 Statewide (Months 7\u201312)",
     "25\u201335 brands. Full NYS coverage with 25+ ambassadors. SMS pipeline across 50+ dispensaries. "
     "Introduce white-label SMS product for dispensaries to purchase directly. "
     "Hire regional managers. Build brand reputation as the go-to sales partner in NY cannabis."),
    ("Phase 4 \u2014 Multi-State (Year 2+)",
     "Replicate the Empire 8 model in NJ, CT, MA, and PA. License the AI marketing and SMS "
     "platform to partner operators. Target: 50+ brands, 50+ ambassadors, 200+ dispensaries on SMS. "
     "Revenue target: $2M+ ARR."),
]

for i, (title, desc) in enumerate(phases):
    y = CT() + Inches(i * 0.75)
    tb(slide, CL(), y, CWF(), Inches(0.2),
       title, size=10, color=PURPLE_ROYAL, bold=True)
    tb(slide, CL(), y + Inches(0.22), CWF(), Inches(0.48),
       desc, size=8, color=TEXT_MED, lsp=1.3)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 22: SCALING ECONOMICS (NEW)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_interior_slide(slide)
section_label(slide, "Scaling Economics")
slide_title(slide, "Why More Brands = Higher Margins")

scaling = [
    "Shared ambassadors \u2014 one rep carries 3\u20135 non-competing brands per territory",
    "Fixed route costs \u2014 same drive whether carrying 1 brand or 5",
    "Content multiplier \u2014 one visit produces content for multiple brands",
    "SMS infrastructure \u2014 same pipeline serves unlimited brands",
    "Budtender relationships compound \u2014 trusted source for multiple products",
    "Data flywheel \u2014 more brands = more data = smarter deployment",
]
bullets(slide, CL(), CT(), CWF(), CH(),
        scaling, size=9, lsp=1.45)

tb(slide, CL(), FRAME_BOTTOM - Inches(0.7), CWF(), Inches(0.3),
   "Linear revenue growth. Sub-linear cost growth. That\u2019s the model.",
   size=10, color=GOLD, bold=True, italic=True)

footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 23: CONTACT / CTA
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PURPLE_DARK)

slide.shapes.add_picture(HEADER_IMG, Emu(0), Emu(0), SLIDE_W, Inches(4.0))

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Inches(3.2), SLIDE_W, Inches(4.3))
shape.fill.solid()
shape.fill.fore_color.rgb = PURPLE_DARK
shape.line.fill.background()

logo_h2 = Inches(4.5)
logo_w2 = Inches(4.5 * (853.0 / 1280.0))
logo_x2 = SLIDE_W - logo_w2 - Inches(0.8)
slide.shapes.add_picture(LOGO_IMG, logo_x2, Inches(0.3), logo_w2, logo_h2)

tb(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
   "READY TO MOVE PRODUCT?",
   size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

tb(slide, Inches(2.5), Inches(4.5), Inches(8.3), Inches(0.8),
   "Schedule a 30-minute meeting with our team.\n"
   "We\u2019ll show you exactly how Empire 8 can grow your brand\u2019s\n"
   "presence in dispensaries across New York.",
   size=14, color=WHITE_70, align=PP_ALIGN.CENTER, lsp=1.6)

tb(slide, Inches(2), Inches(5.6), Inches(9.3), Inches(0.4),
   "mikesutherland@empire8.com",
   size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

tb(slide, Inches(2), Inches(6.1), Inches(9.3), Inches(0.4),
   "empire8ny.com  \u2022  Licensed OCM Distributor  \u2022  New York State",
   size=11, color=WHITE_50, align=PP_ALIGN.CENTER)


# ── Save ──
prs.save(OUT_PATH)
print(f"Deck saved to: {OUT_PATH}")
print(f"{len(prs.slides)} slides generated")
